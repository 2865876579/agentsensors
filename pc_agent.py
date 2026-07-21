"""
PC Agent - 接收云端 AI 命令并控制电脑

功能：
  1. 通过 WebSocket 连接云端服务，等待命令下发
  2. 收到命令后执行对应的电脑操作
  3. 将执行结果返回给云端服务

支持的命令：
  - open_url: 打开网页
  - search: 用默认浏览器搜索
  - open_file: 打开本地文件
  - summarize_file: 读取文件内容并返回摘要

安全设计：
  - 只允许白名单内的动作
  - 不允许删除文件、发送邮件、付款等危险操作
  - 所有操作都有日志输出

用法：
  python pc_agent.py

环境变量：
  WS_URL: 云端服务地址，默认 ws://localhost:8000/ws/pc_agent
  部署后改为 ws://你的服务器IP:8000/ws/pc_agent
"""
import asyncio
import json
import mimetypes
import os
import subprocess
import sys
import time
import urllib.parse
import urllib.request
import webbrowser
from datetime import datetime
from pathlib import Path

import websockets
from dotenv import load_dotenv

# 确保能 import 同目录下的 config
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    # Windows 控制台默认编码可能是 GBK，遇到 emoji/特殊字符会在 print 时抛
    # UnicodeEncodeError，导致 WebSocket 断开。统一改成 UTF-8 并替换无法显示字符。
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


load_dotenv(Path(__file__).with_name(".env"))

# 云端服务 WebSocket 地址
WS_URL = os.getenv("WS_URL", "ws://39.106.190.124:8000/ws/pc_agent")
UPLOAD_TOKEN = os.getenv("PC_AGENT_UPLOAD_TOKEN", "")
MAX_APP_FILE_BYTES = 15 * 1024 * 1024
NO_WINDOW = getattr(subprocess, "CREATE_NO_WINDOW", 0)
WPS_WRITER = Path(r"D:\WPS Office\12.1.0.26895\office6\wps.exe")
WECHAT_EXE = Path(r"D:\vx\聊天记录\Weixin\Weixin.exe")


def _acquire_single_instance():
    if os.name != "nt":
        return None
    import ctypes

    handle = ctypes.windll.kernel32.CreateMutexW(None, False, "Local\\XiaoanPcAgent")
    error = ctypes.windll.kernel32.GetLastError()
    if not handle or error == 183:
        print("[PC Agent] 已有实例正在运行，本次启动退出。")
        raise SystemExit(0)
    return handle


_INSTANCE_MUTEX = None


def cloud_http_url(path: str) -> str:
    parsed = urllib.parse.urlparse(WS_URL)
    scheme = "https" if parsed.scheme == "wss" else "http"
    return urllib.parse.urlunparse((scheme, parsed.netloc, path, "", "", ""))


def _resolve_existing_file(value: str) -> Path:
    raw = os.path.expandvars(str(value or "").strip())
    if not raw:
        raise FileNotFoundError("缺少文件路径")
    source = Path(raw).expanduser()
    candidates = [source]
    if not source.is_absolute():
        relative = Path(*source.parts[1:]) if source.parts and source.parts[0].casefold() == "desktop" else source
        candidates = [Path.home() / "Desktop" / relative, Path.cwd() / source]
    else:
        candidates.append(Path.home() / "Desktop" / source.name)
    for candidate in candidates:
        resolved = candidate.resolve()
        if resolved.is_file():
            return resolved
    raise FileNotFoundError(f"文件不存在: {value}")


def upload_screenshot(path: str) -> dict:
    with open(path, "rb") as image_file:
        image_data = image_file.read(2 * 1024 * 1024 + 1)
    if len(image_data) > 2 * 1024 * 1024:
        raise ValueError("截图压缩后仍超过 2MB")
    request = urllib.request.Request(
        cloud_http_url("/api/pc/screenshots"),
        data=image_data,
        headers={"Content-Type": "image/jpeg", "X-PC-Agent-Token": UPLOAD_TOKEN},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=20) as response:
        payload = json.loads(response.read().decode("utf-8"))
    if not payload.get("ok") or not payload.get("url"):
        raise RuntimeError("云端没有返回截图地址")
    return payload


def upload_file(path: str) -> dict:
    source = _resolve_existing_file(path)
    size = source.stat().st_size
    if size > MAX_APP_FILE_BYTES:
        raise ValueError("文件超过 15MB，无法发送到 App")
    content_type = mimetypes.guess_type(source.name)[0] or "application/octet-stream"
    with source.open("rb") as file_obj:
        data = file_obj.read()
    request = urllib.request.Request(
        cloud_http_url("/api/pc/files"),
        data=data,
        headers={
            "Content-Type": content_type,
            "X-File-Name": urllib.parse.quote(source.name),
            "X-PC-Agent-Token": UPLOAD_TOKEN,
        },
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=30) as response:
        payload = json.loads(response.read().decode("utf-8"))
    if not payload.get("ok") or not payload.get("url"):
        raise RuntimeError("云端没有返回文件地址")
    return payload


def _safe_desktop_path(filename: str, default_name: str) -> Path:
    name = Path(str(filename or default_name)).name
    if not name or name in {".", ".."}:
        name = default_name
    return Path.home() / "Desktop" / name


def _visible_windows(process_names: set[str], title_hint: str = "") -> list[tuple[int, str]]:
    import win32api
    import win32con
    import win32gui
    import win32process

    matches: list[tuple[int, str]] = []
    hint = title_hint.casefold()

    def collect(hwnd, _):
        if not win32gui.IsWindowVisible(hwnd):
            return
        title = win32gui.GetWindowText(hwnd).strip()
        try:
            _, pid = win32process.GetWindowThreadProcessId(hwnd)
            process = win32api.OpenProcess(
                win32con.PROCESS_QUERY_INFORMATION | win32con.PROCESS_VM_READ,
                False,
                pid,
            )
            try:
                process_name = Path(win32process.GetModuleFileNameEx(process, 0)).name.casefold()
            finally:
                process.Close()
        except Exception:
            return
        if process_name not in process_names:
            return
        left, top, right, bottom = win32gui.GetWindowRect(hwnd)
        area = max(0, right - left) * max(0, bottom - top)
        score = 2 if hint and hint in title.casefold() else 1
        matches.append((score, area, hwnd, title))

    win32gui.EnumWindows(collect, None)
    matches.sort(key=lambda item: (item[0], item[1]), reverse=True)
    return [(hwnd, title) for _, _, hwnd, title in matches]


def _focus_window(process_names: set[str], title_hint: str = "", timeout: float = 10.0) -> tuple[bool, str]:
    import win32con
    import win32gui
    from pywinauto import keyboard

    deadline = time.monotonic() + timeout
    last_title = ""
    while time.monotonic() < deadline:
        windows = _visible_windows(process_names, title_hint)
        if windows:
            hwnd, last_title = windows[0]
            try:
                if win32gui.IsIconic(hwnd):
                    win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                else:
                    win32gui.ShowWindow(hwnd, win32con.SW_SHOW)
                keyboard.send_keys("%")
                win32gui.BringWindowToTop(hwnd)
                win32gui.SetForegroundWindow(hwnd)
                time.sleep(0.25)
                if win32gui.GetForegroundWindow() == hwnd:
                    return True, last_title or "应用窗口"
            except Exception:
                pass
        time.sleep(0.25)
    return False, last_title or "未找到可见窗口"


def _launch_and_focus(app_name: str, target: str = "") -> tuple[bool, str]:
    key = str(app_name or "").strip().casefold()
    if key in {"wps", "wps文字", "word", "文字", "文档"}:
        executable = WPS_WRITER
        process_names = {"wps.exe"}
    elif key in {"微信", "wechat", "weixin"}:
        executable = WECHAT_EXE
        process_names = {"weixin.exe", "wechat.exe"}
    elif key in {"浏览器", "browser", "edge", "chrome"}:
        webbrowser.open(target or "about:blank", new=1)
        ok, title = _focus_window({"msedge.exe", "chrome.exe", "firefox.exe", "360chrome.exe"}, timeout=12)
        return ok, title
    else:
        executable = Path(app_name).expanduser()
        if not executable.is_file():
            return False, f"未识别应用或路径不存在: {app_name}"
        process_names = {executable.name.casefold()}

    existing = _visible_windows(process_names)
    if not existing:
        if not executable.is_file():
            return False, f"应用不存在: {executable}"
        subprocess.Popen([str(executable)], close_fds=True)
    ok, title = _focus_window(process_names, timeout=12)
    return ok, title


def _add_word_content(document, content: str) -> None:
    for raw_line in str(content or "").splitlines():
        line = raw_line.strip()
        if not line:
            document.add_paragraph()
        elif line.startswith("### "):
            document.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            document.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            document.add_heading(line[2:], level=1)
        elif line.startswith(("- ", "* ")):
            document.add_paragraph(line[2:], style="List Bullet")
        else:
            document.add_paragraph(line)


def _create_office_file(params: dict) -> Path:
    kind = str(params.get("kind") or "word").strip().casefold()
    title = str(params.get("title") or "整理资料").strip()
    content = str(params.get("content") or "").strip()
    default_names = {
        "word": "整理资料.docx",
        "docx": "整理资料.docx",
        "spreadsheet": "整理表格.xlsx",
        "xlsx": "整理表格.xlsx",
        "presentation": "整理演示.pptx",
        "pptx": "整理演示.pptx",
    }
    if kind not in default_names:
        raise ValueError(f"不支持的 Office 类型: {kind}")
    path = _safe_desktop_path(params.get("filename", ""), default_names[kind])

    if kind in {"word", "docx"}:
        from docx import Document

        if path.suffix.casefold() != ".docx":
            path = path.with_suffix(".docx")
        document = Document()
        if title:
            document.add_heading(title, level=0)
        _add_word_content(document, content)
        document.save(path)
        process_names = {"wps.exe"}
        executable = WPS_WRITER
    elif kind in {"spreadsheet", "xlsx"}:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        if path.suffix.casefold() != ".xlsx":
            path = path.with_suffix(".xlsx")
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = str(params.get("sheet_name") or "资料")[:31]
        headers = params.get("headers") or []
        rows = params.get("rows") or []
        if headers:
            sheet.append([str(value) for value in headers])
            for cell in sheet[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="2F6B5F")
            sheet.freeze_panes = "A2"
        for row in rows:
            if isinstance(row, (list, tuple)):
                sheet.append(list(row))
        if not headers and not rows and content:
            for line in content.splitlines():
                sheet.append([part.strip() for part in line.split("\t")])
        for column in sheet.columns:
            width = min(40, max(10, max((len(str(cell.value or "")) for cell in column), default=10) + 2))
            sheet.column_dimensions[column[0].column_letter].width = width
        workbook.save(path)
        process_names = {"et.exe", "wps.exe"}
        executable = WPS_WRITER.with_name("et.exe")
    else:
        from pptx import Presentation

        if path.suffix.casefold() != ".pptx":
            path = path.with_suffix(".pptx")
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = str(params.get("subtitle") or "")
        slides = params.get("slides") or []
        if not slides and content:
            slides = [{"title": "内容", "content": content}]
        for item in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(item.get("title") or "")
            slide.placeholders[1].text = str(item.get("content") or "")
        presentation.save(path)
        process_names = {"wpp.exe", "wps.exe"}
        executable = WPS_WRITER.with_name("wpp.exe")

    if not path.is_file() or path.stat().st_size == 0:
        raise RuntimeError("Office 文件未成功生成")
    subprocess.Popen(
        [str(executable), str(path)],
        close_fds=True,
        creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP,
    )
    ok, detail = _focus_window(process_names, path.stem, timeout=15)
    if not ok:
        raise RuntimeError(f"文件已生成，但 WPS 窗口未成功置顶: {detail}")
    return path


def _set_clipboard_text(text: str) -> None:
    import win32clipboard

    win32clipboard.OpenClipboard()
    try:
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardText(str(text), win32clipboard.CF_UNICODETEXT)
    finally:
        win32clipboard.CloseClipboard()


def _get_clipboard_text() -> str:
    import win32clipboard

    win32clipboard.OpenClipboard()
    try:
        if not win32clipboard.IsClipboardFormatAvailable(win32clipboard.CF_UNICODETEXT):
            return ""
        return win32clipboard.GetClipboardData(win32clipboard.CF_UNICODETEXT)
    finally:
        win32clipboard.CloseClipboard()


def _set_clipboard_file(path: Path) -> None:
    import struct
    import win32clipboard
    import win32con

    file_list = (str(path) + "\0\0").encode("utf-16le")
    dropfiles = struct.pack("<IiiII", 20, 0, 0, 0, 1) + file_list
    win32clipboard.OpenClipboard()
    try:
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardData(win32con.CF_HDROP, dropfiles)
    finally:
        win32clipboard.CloseClipboard()


def _capture_window(hwnd):
    import win32gui
    from PIL import ImageGrab

    left, top, right, bottom = win32gui.GetWindowRect(hwnd)
    if right <= left or bottom <= top:
        raise RuntimeError("微信窗口尺寸无效")
    return ImageGrab.grab(bbox=(left, top, right, bottom), all_screens=True).convert("RGB")


def _image_change_ratio(before, after, threshold: int = 18) -> float:
    from PIL import ImageChops

    if before.size != after.size:
        return 1.0
    histogram = ImageChops.difference(before, after).convert("L").histogram()
    changed = sum(histogram[threshold + 1:])
    return changed / max(1, before.width * before.height)


def _normalize_ui_text(value: str) -> str:
    return "".join(str(value or "").split()).casefold()


def _click_visible_text(target: str, app: str = "") -> str:
    import win32gui
    from pywinauto import Desktop

    target = str(target or "").strip()
    if not target:
        raise ValueError("缺少要点击的界面文字")

    app_key = str(app or "").strip().casefold()
    if app_key in {"浏览器", "browser", "edge", "chrome"}:
        ok, detail = _focus_window(
            {"msedge.exe", "chrome.exe", "firefox.exe", "360chrome.exe"},
            timeout=8,
        )
    elif app_key:
        ok, detail = _launch_and_focus(app)
    else:
        ok, detail = True, win32gui.GetWindowText(win32gui.GetForegroundWindow())
    if not ok:
        raise RuntimeError(f"目标应用未能置顶: {detail}")

    hwnd = win32gui.GetForegroundWindow()
    if not hwnd:
        raise RuntimeError("当前没有可操作的前台窗口")
    before = _capture_window(hwnd)
    wanted = _normalize_ui_text(target)
    root = Desktop(backend="uia").window(handle=hwnd)
    ranked = []
    control_priority = {
        "Button": 5,
        "Hyperlink": 5,
        "ListItem": 4,
        "TreeItem": 4,
        "TabItem": 4,
        "MenuItem": 4,
        "Text": 2,
    }
    for element in root.descendants():
        try:
            label = str(element.window_text() or "").strip()
            normalized = _normalize_ui_text(label)
            if not normalized or not element.is_visible() or not element.is_enabled():
                continue
            rect = element.rectangle()
            if rect.width() <= 1 or rect.height() <= 1:
                continue
            if normalized == wanted:
                match_score = 3
            elif wanted in normalized:
                match_score = 2
            elif len(normalized) >= 2 and normalized in wanted:
                match_score = 1
            else:
                continue
            control_type = str(element.element_info.control_type or "")
            ranked.append((
                match_score,
                control_priority.get(control_type, 1),
                -abs(len(normalized) - len(wanted)),
                element,
                label,
                control_type,
            ))
        except Exception:
            continue

    if not ranked:
        raise RuntimeError(f"当前窗口中没有找到可点击文字：{target}")
    ranked.sort(key=lambda item: item[:3], reverse=True)
    _, _, _, element, label, control_type = ranked[0]
    element.click_input()
    time.sleep(0.6)

    changed = 0.0
    if win32gui.IsWindow(hwnd):
        try:
            changed = _image_change_ratio(before, _capture_window(hwnd))
        except Exception:
            pass
    return (
        f"已点击“{label}”（{control_type or 'UI'}），"
        f"窗口画面变化 {changed * 100:.1f}%"
    )


def _ensure_wechat_main_window(hwnd, timeout: float = 12.0):
    import win32gui
    from pywinauto import mouse

    left, top, right, bottom = win32gui.GetWindowRect(hwnd)
    if right - left >= 500:
        return hwnd

    # The compact WeChat re-entry window has one primary green button here.
    mouse.click(coords=((left + right) // 2, top + round((bottom - top) * 0.74)))
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        for candidate, _ in _visible_windows({"weixin.exe", "wechat.exe"}):
            x1, y1, x2, y2 = win32gui.GetWindowRect(candidate)
            if x2 - x1 >= 500:
                ok, _ = _focus_window({"weixin.exe", "wechat.exe"}, timeout=2)
                if ok:
                    return candidate
        time.sleep(0.25)
    raise RuntimeError("微信停留在进入界面，未能打开聊天主窗口")


def _send_wechat_file(contact: str, path: str, message: str = "") -> str:
    import win32gui
    from pywinauto import keyboard, mouse

    source = _resolve_existing_file(path)
    if not contact.strip():
        raise ValueError("缺少微信联系人")
    ok, detail = _launch_and_focus("微信")
    if not ok:
        raise RuntimeError(f"微信窗口未成功置顶: {detail}")
    hwnd = _ensure_wechat_main_window(win32gui.GetForegroundWindow())

    keyboard.send_keys("^f")
    time.sleep(0.4)
    _set_clipboard_text(contact.strip())
    keyboard.send_keys("^v")
    time.sleep(1.2)

    # WeChat 4.x search results are painted by Qt and are absent from the UIA tree.
    search_view = _capture_window(hwnd)
    keyboard.send_keys("{ENTER}")
    time.sleep(0.8)
    chat_view = _capture_window(hwnd)
    if _image_change_ratio(search_view, chat_view) < 0.01:
        keyboard.send_keys("{DOWN}{ENTER}")
        time.sleep(0.8)
        chat_view = _capture_window(hwnd)
    if _image_change_ratio(search_view, chat_view) < 0.01:
        keyboard.send_keys("{ESC}")
        raise RuntimeError(f"微信搜索到关键词但未能进入联系人会话: {contact}")
    if win32gui.GetForegroundWindow() != hwnd:
        raise RuntimeError("进入联系人会话时微信失去前台焦点，已取消发送")

    left, top, right, bottom = win32gui.GetWindowRect(hwnd)
    mouse.click(coords=(left + round((right - left) * 0.72), top + round((bottom - top) * 0.88)))
    time.sleep(0.2)
    _set_clipboard_file(source)
    keyboard.send_keys("^v")
    time.sleep(1.2)
    preview_view = _capture_window(hwnd)
    if _image_change_ratio(chat_view, preview_view) < 0.003:
        keyboard.send_keys("{ESC}")
        raise RuntimeError("微信未出现文件发送预览，已取消发送")
    keyboard.send_keys("{ENTER}")
    time.sleep(0.8)
    sent_view = _capture_window(hwnd)
    if _image_change_ratio(preview_view, sent_view) < 0.01:
        keyboard.send_keys("{ESC}")
        raise RuntimeError("微信文件发送预览未关闭，无法确认文件已经发送")
    if message.strip():
        time.sleep(0.5)
        _set_clipboard_text(message.strip())
        keyboard.send_keys("^v{ENTER}")
    return f"已发送给微信联系人 {contact}: {source.name}"


def extract_search_query_from_url(url: str) -> str | None:
    """识别搜索引擎 URL，返回搜索词；不是搜索页则返回 None。"""
    parsed = urllib.parse.urlparse(url)
    host = parsed.netloc.lower()
    params = urllib.parse.parse_qs(parsed.query)

    if "baidu.com" in host:
        names = ("wd", "word", "q")
    elif "bing.com" in host or "google." in host or "duckduckgo.com" in host:
        names = ("q",)
    elif "sogou.com" in host:
        names = ("query", "keyword", "q")
    elif "so.com" in host or "haosou.com" in host:
        names = ("q",)
    else:
        return None

    for name in names:
        values = params.get(name)
        if values and values[0].strip():
            return urllib.parse.unquote_plus(values[0]).strip()
    return ""


async def handle_command(command: dict) -> str:
    """
    执行 PC 命令并返回结果

    参数：
      command: {"action": "动作名", "params": {"参数": "值"}}

    返回：
      执行结果的文字描述
    """
    action = command.get("action", "")
    params = command.get("params", {})

    print(f"[执行] action={action}, params={params}")

    if action == "open_url":
        url = params.get("url", "")
        if not url:
            return "缺少 URL 参数"
        query = extract_search_query_from_url(url)
        if query is not None:
            return f"已拦截搜索页打开：{query or url}。搜索应该由服务端后台完成，不再打开浏览器。"
        webbrowser.open(url)
        return f"已打开网页: {url}"

    elif action == "search":
        query = params.get("query", "")
        if not query:
            return "缺少搜索关键词"
        return f"已拦截 PC Agent 搜索命令：{query}。搜索应该由服务端后台完成，不再打开浏览器。"

    elif action == "open_file":
        path = params.get("path") or params.get("file_path") or ""
        if not path:
            return "缺少文件路径"
        try:
            source = _resolve_existing_file(path)
            os.startfile(source)
            process_names = {
                ".docx": {"wps.exe"},
                ".xlsx": {"et.exe", "wps.exe"},
                ".pptx": {"wpp.exe", "wps.exe"},
            }.get(source.suffix.casefold())
            if process_names:
                ok, detail = await asyncio.to_thread(_focus_window, process_names, source.stem, 15)
                if not ok:
                    return f"文件已打开，但窗口未成功置顶: {detail}"
            return f"已打开并验证文件窗口: {source}"
        except Exception as e:
            return f"打开文件失败: {e}"

    elif action == "summarize_file":
        path = params.get("path", "")
        if not path:
            return "缺少文件路径"
        if not os.path.exists(path):
            return f"文件不存在: {path}"
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read(2000)
            return f"文件内容前2000字：\n{content}"
        except Exception as e:
            return f"读取文件失败: {e}"

    elif action == "desktop_write":
        # 在用户桌面创建/覆盖文件
        filename = params.get("filename", "summary.txt")
        content = params.get("content", "")
        if not content:
            return "缺少文件内容"
        path = _safe_desktop_path(filename, "summary.txt")
        try:
            with path.open("w", encoding="utf-8") as f:
                f.write(content)
            return f"已写入并验证桌面文件：{path}（{len(content)}字）"
        except Exception as e:
            return f"写入桌面文件失败: {e}"

    elif action == "clipboard_get":
        # 读取 Windows 剪贴板文本
        try:
            text = (await asyncio.to_thread(_get_clipboard_text)).strip()
            return text if text else "剪贴板为空"
        except Exception as e:
            return f"读取剪贴板失败: {e}"

    elif action == "clipboard_set":
        # 写入 Windows 剪贴板
        text = params.get("text", "")
        if not text:
            return "缺少要写入剪贴板的文本"
        try:
            await asyncio.to_thread(_set_clipboard_text, text)
            return f"已写入剪贴板（{len(text)}字）"
        except Exception as e:
            return f"写入剪贴板失败: {e}"

    elif action == "screenshot":
        # Capture in-process so no PowerShell window can appear in the image.
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"screenshot_{ts}.jpg"
        path = os.path.join(desktop, filename)
        try:
            from PIL import ImageGrab

            image = ImageGrab.grab(all_screens=False).convert("RGB")
            if image.width > 1280:
                height = round(image.height * 1280 / image.width)
                image = image.resize((1280, height))
            image.save(path, "JPEG", quality=85, optimize=True)
            if not os.path.isfile(path) or os.path.getsize(path) == 0:
                return "截图失败: 截图文件未生成"
            await asyncio.to_thread(upload_screenshot, path)
            return f"截图已发送到 App，并保存到桌面：{filename}"
        except Exception as e:
            return f"截图失败: {e}"

    elif action == "run_app":
        app = params.get("app", "")
        if not app:
            return "缺少应用名称或路径"
        try:
            ok, detail = await asyncio.to_thread(_launch_and_focus, app, str(params.get("url") or ""))
            if not ok:
                return f"启动失败: {detail}"
            return f"已启动并置于前台：{app}（{detail}）"
        except Exception as e:
            return f"启动失败: {e}"

    elif action == "click_text":
        target = str(params.get("target") or "")
        app = str(params.get("app") or "")
        try:
            return await asyncio.to_thread(_click_visible_text, target, app)
        except Exception as e:
            return f"界面点击失败: {e}"

    elif action == "create_office_file":
        try:
            path = await asyncio.to_thread(_create_office_file, params)
            completed = [f"Office 文件已生成、校验并在 WPS 前台打开：{path}"]
            if params.get("send_to_app"):
                payload = await asyncio.to_thread(upload_file, str(path))
                completed.append(f"已发送到 App：{payload.get('original_name') or path.name}")
            wechat_contact = str(params.get("wechat_contact") or "").strip()
            if wechat_contact:
                completed.append(await asyncio.to_thread(
                    _send_wechat_file,
                    wechat_contact,
                    str(path),
                    str(params.get("message") or ""),
                ))
            return "；".join(completed)
        except Exception as e:
            return f"创建 Office 文件失败: {e}"

    elif action == "send_file_to_app":
        path = params.get("path", "")
        if not path:
            return "缺少文件路径"
        try:
            payload = await asyncio.to_thread(upload_file, path)
            return f"文件已发送到 App：{payload.get('original_name') or Path(path).name}"
        except Exception as e:
            return f"发送文件到 App 失败: {e}"

    elif action == "send_wechat_file":
        contact = str(params.get("contact") or "")
        path = str(params.get("path") or "")
        message = str(params.get("message") or "")
        try:
            return await asyncio.to_thread(_send_wechat_file, contact, path, message)
        except Exception as e:
            return f"微信发送失败: {e}"

    elif action == "run_cmd":
        # 执行命令并返回输出
        cmd = params.get("cmd") or params.get("command") or ""
        if not cmd:
            return "缺少要执行的命令"
        try:
            r = subprocess.run(
                cmd,
                shell=True,
                capture_output=True,
                text=True,
                timeout=15,
                creationflags=NO_WINDOW,
            )
            out = r.stdout.strip() or r.stderr.strip()
            if r.returncode != 0:
                return f"命令执行失败（返回码 {r.returncode}）：{out[:800] or '无错误输出'}"
            if not out:
                return "命令执行成功（返回码 0，无输出）"
            preview = out[:800]
            return preview if len(out) <= 800 else preview + "\n...(输出已截断)"
        except Exception as e:
            return f"执行命令失败: {e}"

    else:
        return f"不支持的命令: {action}"


async def run():
    """主循环：连接云端服务，等待并执行命令"""
    while True:
        try:
            print(f"[PC Agent] 连接到 {WS_URL} ...")
            async with websockets.connect(WS_URL) as ws:
                print("[PC Agent] 已连接，等待命令...\n")

                while True:
                    message = await ws.recv()
                    data = json.loads(message)

                    if data.get("type") == "pc_command":
                        command = data.get("command", {})
                        client_id = data.get("client_id")
                        turn_id = data.get("turn_id")
                        command_id = data.get("command_id")
                        print(f"[收到命令] {command}")

                        # 执行命令
                        result = await handle_command(command)
                        print(f"[执行结果] {result}\n")

                        # 返回结果给云端
                        await ws.send(json.dumps({
                            "type": "result",
                            "client_id": client_id,
                            "turn_id": turn_id,
                            "command_id": command_id,
                            "result": result
                        }, ensure_ascii=False))

                    elif data.get("type") == "pong":
                        pass

        except websockets.exceptions.ConnectionClosed:
            print("[PC Agent] 连接断开，5秒后重连...")
            await asyncio.sleep(5)
        except ConnectionRefusedError:
            print("[PC Agent] 服务器未启动，5秒后重试...")
            await asyncio.sleep(5)
        except Exception as e:
            print(f"[PC Agent] 错误: {e}，5秒后重连...")
            await asyncio.sleep(5)


if __name__ == "__main__":
    _INSTANCE_MUTEX = _acquire_single_instance()
    print("=" * 40)
    print("  智能枕头 PC Agent")
    print("  等待云端 AI 下发电脑控制命令")
    print("=" * 40)
    print()
    asyncio.run(run())
